import os
import io
import time
import uuid
import hashlib
from typing import List, Dict, Any, Tuple, Optional

from fastapi import UploadFile

# Text extraction libs
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
from PIL import Image

# Optional OCR for scanned PDFs
import pytesseract
try:
    from pdf2image import convert_from_path  # needs poppler on the system/Docker
    PDF2IMAGE_AVAILABLE = True
except Exception:
    PDF2IMAGE_AVAILABLE = False

from .indexing import get_chroma_collection

# ----------------------------
# Config (chunking + OCR)
# ----------------------------
CHUNK_SIZE = int(os.getenv("CHUNK_SIZE", "1000"))       # target chars per chunk
CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", "200"))  # overlap between chunks
OCR_ENABLED = os.getenv("OCR_ENABLED", "true").lower() in ("1", "true", "yes")


# ----------------------------
# Helpers
# ----------------------------
def _sha256_file(path: str) -> str:
    """Compute SHA-256 hash of a file."""
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            h.update(chunk)
    return h.hexdigest()


def _ensure_dirs() -> Tuple[str, str]:
    """Ensure data directories exist; return (uploads_dir, processed_dir)."""
    uploads = os.path.join("data", "uploads")
    processed = os.path.join("data", "processed")
    os.makedirs(uploads, exist_ok=True)
    os.makedirs(processed, exist_ok=True)
    return uploads, processed


def _save_upload(file: UploadFile) -> Tuple[str, str]:
    """Save UploadFile to disk and return (saved_path, filename)."""
    uploads_dir, _ = _ensure_dirs()
    filename = file.filename or f"upload_{uuid.uuid4().hex}"
    # Normalize path separators for Windows/Unix
    safe_name = filename.replace("/", "_").replace("\\", "_")
    saved_path = os.path.join(uploads_dir, safe_name)
    with open(saved_path, "wb") as out:
        out.write(file.file.read())
    return saved_path, safe_name


def _text_from_pdf_with_pdfplumber(path: str) -> str:
    text_parts: List[str] = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            try:
                txt = page.extract_text() or ""
                if txt.strip():
                    text_parts.append(txt)
            except Exception:
                # Skip problematic pages
                continue
    return "\n".join(text_parts).strip()


def _ocr_pdf(path: str, dpi: int = 200) -> str:
    """OCR each page of a PDF (if pdf2image/poppler is available)."""
    if not PDF2IMAGE_AVAILABLE:
        return ""
    try:
        images = convert_from_path(path, dpi=dpi)
    except Exception:
        return ""
    text_parts: List[str] = []
    for img in images:
        try:
            txt = pytesseract.image_to_string(img) or ""
            if txt.strip():
                text_parts.append(txt)
        except Exception:
            continue
    return "\n".join(text_parts).strip()


def _text_from_docx(path: str) -> str:
    doc = DocxDocument(path)
    lines = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    return "\n".join(lines).strip()


def _text_from_pptx(path: str) -> str:
    prs = Presentation(path)
    lines: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                t = shape.text.strip()
                if t:
                    lines.append(t)
    return "\n".join(lines).strip()


def _text_from_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read().strip()


def extract_text_any(path: str, filename: str) -> str:
    """Extract text from PDF/DOCX/PPTX/TXT, with optional OCR for PDFs."""
    lower = filename.lower()
    if lower.endswith(".pdf"):
        # Try native text first
        text = _text_from_pdf_with_pdfplumber(path)
        if not text and OCR_ENABLED:
            # Fallback to OCR for scanned PDFs
            text = _ocr_pdf(path)
        return text or ""
    elif lower.endswith(".docx"):
        return _text_from_docx(path)
    elif lower.endswith(".pptx"):
        return _text_from_pptx(path)
    elif lower.endswith(".txt"):
        return _text_from_txt(path)
    else:
        # Unknown type → best-effort read as text
        try:
            return _text_from_txt(path)
        except Exception:
            return ""


def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    """Simple character-based chunking with overlap."""
    if not text:
        return []
    chunks: List[str] = []
    start = 0
    n = len(text)
    while start < n:
        end = min(n, start + chunk_size)
        chunk = text[start:end]
        chunks.append(chunk)
        if end == n:
            break
        start = max(end - overlap, start + 1)
    return chunks


def _make_chunk_ids(doc_id: str, count: int) -> List[str]:
    return [f"{doc_id}::chunk::{i}" for i in range(count)]


def _make_metadatas(doc_id: str, filename: str, source_path: str, file_hash: str, count: int) -> List[Dict[str, Any]]:
    ts = int(time.time())
    metas: List[Dict[str, Any]] = []
    for i in range(count):
        metas.append({
            "doc_id": doc_id,
            "filename": filename,
            "source_path": source_path,
            "chunk_index": i,
            "ingested_at": ts,
            "file_hash": file_hash,
        })
    return metas


# ----------------------------
# Public API used by FastAPI routes
# ----------------------------
def save_upload(filename: str, raw: bytes) -> str:
    """Save raw bytes to uploads directory and return saved path."""
    uploads_dir, _ = _ensure_dirs()
    # Normalize path separators for Windows/Unix
    safe_name = filename.replace("/", "_").replace("\\", "_")
    saved_path = os.path.join(uploads_dir, safe_name)
    with open(saved_path, "wb") as out:
        out.write(raw)
    return saved_path

def extract_text_from_pdf(raw: bytes) -> str:
    """Extract text from PDF bytes."""
    import tempfile
    with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
        tmp.write(raw)
        tmp_path = tmp.name
    try:
        return extract_text_any(tmp_path, "temp.pdf")
    finally:
        os.unlink(tmp_path)

def extract_text_from_docx(raw: bytes) -> str:
    """Extract text from DOCX bytes."""
    import tempfile
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(raw)
        tmp_path = tmp.name
    try:
        return extract_text_any(tmp_path, "temp.docx")
    finally:
        os.unlink(tmp_path)

def extract_text_from_pptx(raw: bytes) -> str:
    """Extract text from PPTX bytes."""
    import tempfile
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(raw)
        tmp_path = tmp.name
    try:
        return extract_text_any(tmp_path, "temp.pptx")
    finally:
        os.unlink(tmp_path)

def extract_text_from_txt(raw: bytes) -> str:
    """Extract text from TXT bytes."""
    return raw.decode('utf-8', errors='ignore').strip()

def upload_only(file: UploadFile) -> Dict[str, Any]:
    """Just save the file; return path/meta (no indexing)."""
    saved_path, filename = _save_upload(file)
    return {
        "message": "Uploaded",
        "filename": filename,
        "saved_path": saved_path,
        "size_bytes": os.path.getsize(saved_path),
    }


def ingest(file: UploadFile) -> Dict[str, Any]:
    """
    Save upload -> extract text (OCR for PDFs if needed) -> chunk+embed to Chroma.
    Always returns a dict (never None).
    """
    saved_path, filename = _save_upload(file)

    # Duplicate protection by file hash
    file_hash = _sha256_file(saved_path)
    coll = get_chroma_collection()
    try:
        existing = coll.get(where={"file_hash": file_hash}, include=["ids", "metadatas"])
    except Exception:
        existing = None

    if existing and existing.get("ids") and existing["ids"][0]:
        first_meta = (existing.get("metadatas") or [[]])[0]
        first_meta = first_meta[0] if first_meta else {}
        return {
            "message": "Duplicate skipped (already ingested)",
            "filename": filename,
            "file_hash": file_hash,
            "doc_id": (first_meta or {}).get("doc_id"),
            "chunks": len(existing["ids"][0]),
            "saved_path": saved_path,
        }

    # New ingestion
    text = extract_text_any(saved_path, filename)
    if not text.strip():
        return {
            "message": "No text extracted (empty document or OCR failed)",
            "filename": filename,
            "doc_id": None,
            "chunks": 0,
            "saved_path": saved_path,
        }

    pieces = chunk_text(text, CHUNK_SIZE, CHUNK_OVERLAP)
    doc_id = str(uuid.uuid4())
    ids = _make_chunk_ids(doc_id, len(pieces))
    metadatas = _make_metadatas(doc_id, filename, saved_path, file_hash, len(pieces))

    # Upsert into Chroma
    # Prefer upsert if available; fallback to add.
    try:
        if hasattr(coll, "upsert"):
            coll.upsert(documents=pieces, ids=ids, metadatas=metadatas)
        else:
            coll.add(documents=pieces, ids=ids, metadatas=metadatas)
    except Exception as e:
        return {
            "message": f"Indexing error: {e}",
            "filename": filename,
            "doc_id": doc_id,
            "chunks": 0,
            "saved_path": saved_path,
        }

    return {
        "message": "Ingested",
        "doc_id": doc_id,
        "filename": filename,
        "chunks": len(pieces),
        "saved_path": saved_path,
    }
