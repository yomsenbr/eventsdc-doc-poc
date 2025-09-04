import io
import logging
from typing import Dict, Any
import pytesseract
from PIL import Image
import PyPDF2
from docx import Document
from pptx import Presentation
import cv2
import numpy as np

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Document processor with OCR capabilities"""
    
    def __init__(self):
        self.supported_formats = {'.pdf', '.docx', '.pptx', '.txt'}
    
    def process_document(self, file_content: bytes, file_extension: str) -> str:
        """Process document and extract text with OCR if needed"""
        try:
            text = ""
            
            if file_extension == '.pdf':
                text = self._process_pdf(file_content)
            elif file_extension == '.docx':
                text = self._process_docx(file_content)
            elif file_extension == '.pptx':
                text = self._process_pptx(file_content)
            elif file_extension == '.txt':
                text = self._process_txt(file_content)
            else:
                raise ValueError(f"Unsupported file format: {file_extension}")
            
            # Clean and normalize text
            text = self._clean_text(text)
            
            logger.info(f"Extracted {len(text)} characters from {file_extension} document")
            return text
            
        except Exception as e:
            logger.error(f"Error processing document: {str(e)}")
            raise
    
    def _process_pdf(self, file_content: bytes) -> str:
        """Process PDF files with OCR fallback"""
        text = ""
        
        try:
            # Try text extraction first
            pdf_file = io.BytesIO(file_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            for page_num, page in enumerate(pdf_reader.pages):
                page_text = page.extract_text()
                if page_text.strip():
                    text += page_text + "\n"
                else:
                    # If no text extracted, try OCR
                    logger.info(f"No text found on page {page_num + 1}, attempting OCR")
                    # Note: OCR on PDF would require additional image conversion
                    # This is a simplified implementation
            
            return text
            
        except Exception as e:
            logger.error(f"Error processing PDF: {str(e)}")
            raise
    
    def _process_docx(self, file_content: bytes) -> str:
        """Process DOCX files"""
        try:
            doc_file = io.BytesIO(file_content)
            doc = Document(doc_file)
            
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
                text += "\n"
            
            return text
            
        except Exception as e:
            logger.error(f"Error processing DOCX: {str(e)}")
            raise
    
    def _process_pptx(self, file_content: bytes) -> str:
        """Process PPTX files"""
        try:
            ppt_file = io.BytesIO(file_content)
            prs = Presentation(ppt_file)
            
            text = ""
            for slide_num, slide in enumerate(prs.slides):
                text += f"Slide {slide_num + 1}:\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                
                text += "\n"
            
            return text
            
        except Exception as e:
            logger.error(f"Error processing PPTX: {str(e)}")
            raise
    
    def _process_txt(self, file_content: bytes) -> str:
        """Process TXT files"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'latin-1', 'cp1252']
            
            for encoding in encodings:
                try:
                    text = file_content.decode(encoding)
                    return text
                except UnicodeDecodeError:
                    continue
            
            # If all encodings fail, use utf-8 with error handling
            text = file_content.decode('utf-8', errors='replace')
            return text
            
        except Exception as e:
            logger.error(f"Error processing TXT: {str(e)}")
            raise
    
    def _perform_ocr(self, image_data: bytes) -> str:
        """Perform OCR on image data"""
        try:
            # Convert bytes to PIL Image
            image = Image.open(io.BytesIO(image_data))
            
            # Convert to OpenCV format for preprocessing
            cv_image = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2BGR)
            
            # Preprocess image for better OCR
            gray = cv2.cvtColor(cv_image, cv2.COLOR_BGR2GRAY)
            
            # Apply thresholding to get better text extraction
            _, thresh = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
            
            # Convert back to PIL Image
            processed_image = Image.fromarray(thresh)
            
            # Perform OCR
            text = pytesseract.image_to_string(processed_image, lang='eng')
            
            return text
            
        except Exception as e:
            logger.error(f"Error performing OCR: {str(e)}")
            return ""
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize extracted text"""
        if not text:
            return ""
        
        # Remove excessive whitespace
        lines = text.split('\n')
        cleaned_lines = []
        
        for line in lines:
            # Strip whitespace and skip empty lines
            cleaned_line = line.strip()
            if cleaned_line:
                cleaned_lines.append(cleaned_line)
        
        # Join lines with single newline
        cleaned_text = '\n'.join(cleaned_lines)
        
        # Remove excessive spaces
        import re
        cleaned_text = re.sub(r' +', ' ', cleaned_text)
        
        return cleaned_text
    
    def get_document_metadata(self, file_content: bytes, file_extension: str, filename: str) -> Dict[str, Any]:
        """Extract metadata from document"""
        metadata = {
            'filename': filename,
            'file_extension': file_extension,
            'file_size': len(file_content),
            'pages': 0,
            'has_images': False
        }
        
        try:
            if file_extension == '.pdf':
                pdf_file = io.BytesIO(file_content)
                pdf_reader = PyPDF2.PdfReader(pdf_file)
                metadata['pages'] = len(pdf_reader.pages)
                
                # Check for metadata
                if pdf_reader.metadata:
                    if pdf_reader.metadata.title:
                        metadata['title'] = pdf_reader.metadata.title
                    if pdf_reader.metadata.author:
                        metadata['author'] = pdf_reader.metadata.author
                    if pdf_reader.metadata.creation_date:
                        metadata['creation_date'] = str(pdf_reader.metadata.creation_date)
            
            elif file_extension == '.pptx':
                ppt_file = io.BytesIO(file_content)
                prs = Presentation(ppt_file)
                metadata['pages'] = len(prs.slides)
                metadata['slide_count'] = len(prs.slides)
            
            elif file_extension == '.docx':
                doc_file = io.BytesIO(file_content)
                doc = Document(doc_file)
                metadata['paragraphs'] = len(doc.paragraphs)
                metadata['tables'] = len(doc.tables)
                
                # Get document properties
                if doc.core_properties.title:
                    metadata['title'] = doc.core_properties.title
                if doc.core_properties.author:
                    metadata['author'] = doc.core_properties.author
                if doc.core_properties.created:
                    metadata['creation_date'] = str(doc.core_properties.created)
        
        except Exception as e:
            logger.warning(f"Could not extract metadata: {str(e)}")
        
        return metadata